#!/usr/bin/env python3
"""
대한피부건강산업협회 소개 PDF → 영문 PPTX 변환
사용법: python scripts/translate_kshia_pdf.py

파이프라인:
  1) PDF → 페이지별 이미지 (pdf2image + poppler)
  2) 이미지 → 한글 OCR (Gemini Vision API)
  3) 한글 → 영어 번역 (Gemini API)
  4) 영문 PPTX 생성 (python-pptx)
"""

import os
import sys
import json
import time
import base64
from pathlib import Path

PROJECT_ROOT = Path(__file__).parent.parent
PDF_PATH = PROJECT_ROOT / "test 자료" / "TalkFile_대한피부건강산업협회 소개(26.03.01).pdf.pdf"
OUTPUT_DIR = PROJECT_ROOT / "data" / "kshia_translation"
PPTX_OUTPUT = OUTPUT_DIR / "KSHIA_Introduction_EN.pptx"

# .env에서 API 키 로드
def load_env():
    env_path = PROJECT_ROOT / "config" / ".env"
    if env_path.exists():
        for line in env_path.read_text().splitlines():
            line = line.strip()
            if line and not line.startswith("#") and "=" in line:
                key, val = line.split("=", 1)
                os.environ.setdefault(key.strip(), val.strip())

load_env()

import google.generativeai as genai
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

genai.configure(api_key=os.environ["GEMINI_API_KEY"])


# ─── 1단계: PDF → 이미지 ──────────────────────────────────────

def pdf_to_images():
    """PDF를 페이지별 PNG 이미지로 변환"""
    img_dir = OUTPUT_DIR / "pages"
    img_dir.mkdir(parents=True, exist_ok=True)

    print("📄 PDF → 이미지 변환 중...")
    print(f"   파일: {PDF_PATH.name} ({PDF_PATH.stat().st_size // 1024 // 1024}MB)")

    images = convert_from_path(
        str(PDF_PATH),
        dpi=200,            # OCR에 충분한 해상도
        fmt="png",
        output_folder=str(img_dir),
        paths_only=True,
    )

    # 파일명 정리
    img_paths = []
    for i, img_path in enumerate(images):
        new_path = img_dir / f"page_{i+1:02d}.png"
        if str(img_path) != str(new_path):
            os.rename(img_path, new_path)
        img_paths.append(new_path)
        size_kb = new_path.stat().st_size // 1024
        print(f"   ✅ 페이지 {i+1}: {new_path.name} ({size_kb}KB)")

    print(f"   총 {len(img_paths)}페이지 변환 완료")
    return img_paths


# ─── 2단계: Gemini Vision OCR ─────────────────────────────────

def ocr_pages(img_paths):
    """각 이미지에서 한글 텍스트 추출 (Gemini Vision)"""
    print("\n🔍 OCR 텍스트 추출 중 (Gemini Vision)...")

    model = genai.GenerativeModel("gemini-2.5-flash")
    ocr_results = []

    for img_path in img_paths:
        page_num = img_path.stem.split("_")[1]
        print(f"   페이지 {page_num} OCR 중...", end=" ", flush=True)

        # 이미지를 base64로 인코딩
        img_data = img_path.read_bytes()

        prompt = """이 프레젠테이션 슬라이드 이미지에서 모든 텍스트를 추출해주세요.

규칙:
1. 슬라이드의 구조를 유지하며 추출 (제목, 본문, 표, 목록 구분)
2. 텍스트의 계층 구조를 표현 (제목은 ##, 소제목은 ###, 본문은 일반 텍스트)
3. 표가 있으면 마크다운 표로 변환
4. 인물 이름, 직책, 소속 등 고유명사는 정확히 추출
5. 영어가 섞여 있으면 영어 그대로 유지
6. 빈 칸이나 구분선도 구조적으로 표현

JSON 형식으로 응답해주세요:
{
  "page_type": "cover|content|profile|organization|partner",
  "title": "슬라이드 제목",
  "sections": [
    {
      "type": "heading|paragraph|list|table|profile",
      "content": "텍스트 내용"
    }
  ]
}"""

        try:
            response = model.generate_content([
                prompt,
                {"mime_type": "image/png", "data": img_data},
            ])
            text = response.text.strip()

            # JSON 파싱 시도 (마크다운 코드블록 제거)
            if text.startswith("```"):
                text = text.split("\n", 1)[1]
                if text.endswith("```"):
                    text = text.rsplit("```", 1)[0]
            text = text.strip()

            try:
                result = json.loads(text)
            except json.JSONDecodeError:
                result = {"page_type": "content", "title": f"페이지 {page_num}", "raw_text": text, "sections": []}

            result["page_number"] = int(page_num)
            ocr_results.append(result)
            print(f"✅ ({result.get('page_type', '?')})")

        except Exception as e:
            print(f"❌ 오류: {e}")
            ocr_results.append({
                "page_number": int(page_num),
                "page_type": "error",
                "title": f"페이지 {page_num}",
                "error": str(e),
                "sections": [],
            })

        # API 속도 제한 방지
        time.sleep(2)

    # OCR 결과 저장
    ocr_path = OUTPUT_DIR / "ocr_results.json"
    with open(ocr_path, "w", encoding="utf-8") as f:
        json.dump(ocr_results, f, ensure_ascii=False, indent=2)
    print(f"\n   💾 OCR 결과 저장: {ocr_path}")

    return ocr_results


# ─── 3단계: 번역 ──────────────────────────────────────────────

def translate_pages(ocr_results):
    """OCR 결과를 영어로 번역"""
    print("\n🌐 영어 번역 중 (Gemini)...")

    model = genai.GenerativeModel("gemini-2.5-flash")
    translated = []

    for page in ocr_results:
        page_num = page["page_number"]
        print(f"   페이지 {page_num} 번역 중...", end=" ", flush=True)

        # 페이지 내용을 텍스트로 직렬화
        page_text = json.dumps(page, ensure_ascii=False, indent=2)

        prompt = f"""아래는 한국어 프레젠테이션 슬라이드의 OCR 추출 결과입니다.
이것을 영어로 번역해주세요.

번역 규칙:
1. 조직명: "대한피부건강산업협회" → "Korea Skin Health Industry Association (KSHIA)"
2. 인물 이름은 한국식 영문 표기 (예: 박종수 → Jong-Su Park)
3. 대학/기관명은 공식 영문명 사용 (예: 상하이교통대학 → Shanghai Jiao Tong University)
4. 직책은 영어 관용 표현 사용 (예: 이사장 → Chairman, 부이사장 → Vice Chairman)
5. 학술 약어(SCI, SCOPUS, H-INDEX 등)는 그대로 유지
6. 전문 용어는 정확한 영문 의학/과학 용어 사용
7. JSON 구조를 그대로 유지하되 content만 영어로 변환

원본 OCR:
{page_text}

영어 번역된 JSON을 반환해주세요 (구조 동일):"""

        try:
            response = model.generate_content(prompt)
            text = response.text.strip()

            if text.startswith("```"):
                text = text.split("\n", 1)[1]
                if text.endswith("```"):
                    text = text.rsplit("```", 1)[0]
            text = text.strip()

            try:
                result = json.loads(text)
            except json.JSONDecodeError:
                result = {**page, "translated_raw": text}

            result["page_number"] = page_num
            translated.append(result)
            print(f"✅")

        except Exception as e:
            print(f"❌ 오류: {e}")
            translated.append({**page, "translation_error": str(e)})

        time.sleep(2)

    # 번역 결과 저장
    tr_path = OUTPUT_DIR / "translated_results.json"
    with open(tr_path, "w", encoding="utf-8") as f:
        json.dump(translated, f, ensure_ascii=False, indent=2)
    print(f"\n   💾 번역 결과 저장: {tr_path}")

    return translated


# ─── 4단계: PPTX 생성 ─────────────────────────────────────────

def create_pptx(translated, img_paths):
    """번역된 내용으로 영문 PPTX 생성"""
    print("\n📊 영문 PPTX 생성 중...")

    prs = Presentation()

    # 슬라이드 크기: 와이드스크린 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 색상 팔레트
    NAVY = RGBColor(0x1F, 0x4E, 0x79)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    DARK = RGBColor(0x33, 0x33, 0x33)
    ACCENT = RGBColor(0x2E, 0x75, 0xB6)
    LIGHT_BG = RGBColor(0xF2, 0xF2, 0xF2)

    for i, (page, img_path) in enumerate(zip(translated, img_paths)):
        page_type = page.get("page_type", "content")
        title = page.get("title", f"Page {i+1}")
        sections = page.get("sections", [])

        # 빈 슬라이드 레이아웃 (6 = Blank)
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # 배경색 설정
        background = slide.background
        fill = background.fill
        fill.solid()
        if page_type == "cover":
            fill.fore_color.rgb = NAVY
        else:
            fill.fore_color.rgb = WHITE

        # ── 제목 ──
        title_color = WHITE if page_type == "cover" else NAVY
        title_size = Pt(36) if page_type == "cover" else Pt(28)
        title_top = Inches(2.5) if page_type == "cover" else Inches(0.3)

        title_box = slide.shapes.add_textbox(
            Inches(0.8), title_top, Inches(11.7), Inches(1.2)
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = title_size
        title_para.font.bold = True
        title_para.font.color.rgb = title_color
        title_para.alignment = PP_ALIGN.CENTER if page_type == "cover" else PP_ALIGN.LEFT

        # 표지: 부제목 추가
        if page_type == "cover":
            sub_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(4.0), Inches(11.7), Inches(0.8)
            )
            sub_frame = sub_box.text_frame
            sub_para = sub_frame.paragraphs[0]
            sub_para.text = "KOREA SKIN HEALTH INDUSTRY ASSOCIATION"
            sub_para.font.size = Pt(18)
            sub_para.font.color.rgb = RGBColor(0xBB, 0xCC, 0xDD)
            sub_para.alignment = PP_ALIGN.CENTER
            continue

        # ── 본문 영역 ──
        body_top = Inches(1.5)
        body_box = slide.shapes.add_textbox(
            Inches(0.8), body_top, Inches(11.7), Inches(5.5)
        )
        body_frame = body_box.text_frame
        body_frame.word_wrap = True

        for j, section in enumerate(sections):
            sec_type = section.get("type", "paragraph")
            content = section.get("content", "")

            # content가 dict/list면 문자열로 변환
            if isinstance(content, dict):
                content = json.dumps(content, ensure_ascii=False)
            elif isinstance(content, list) and sec_type != "list":
                content = "\n".join(str(x) for x in content)
            elif not isinstance(content, str):
                content = str(content)

            if not content:
                continue

            if j > 0:
                # 섹션 간 빈 줄
                spacer = body_frame.add_paragraph()
                spacer.space_after = Pt(4)

            if sec_type == "heading":
                para = body_frame.add_paragraph() if j > 0 else body_frame.paragraphs[0]
                para.text = content
                para.font.size = Pt(20)
                para.font.bold = True
                para.font.color.rgb = ACCENT
                para.space_after = Pt(8)

            elif sec_type == "list":
                # 리스트 항목 (줄바꿈으로 분리 또는 배열)
                if isinstance(content, list):
                    items = content
                elif "\n" in content:
                    items = content.split("\n")
                else:
                    items = [content]
                for item in items:
                    if isinstance(item, (dict, list)):
                        item = json.dumps(item, ensure_ascii=False) if isinstance(item, dict) else str(item)
                    item = str(item).strip().lstrip("•·-▪►● ")
                    if not item:
                        continue
                    para = body_frame.add_paragraph() if j > 0 or body_frame.paragraphs[0].text else body_frame.paragraphs[0]
                    para.text = f"  •  {item}"
                    para.font.size = Pt(14)
                    para.font.color.rgb = DARK
                    para.space_after = Pt(4)

            elif sec_type == "table":
                # 표는 텍스트로 대체 (간략 처리)
                para = body_frame.add_paragraph() if j > 0 else body_frame.paragraphs[0]
                para.text = content
                para.font.size = Pt(12)
                para.font.color.rgb = DARK
                para.space_after = Pt(6)

            elif sec_type == "profile":
                para = body_frame.add_paragraph() if j > 0 else body_frame.paragraphs[0]
                para.text = content
                para.font.size = Pt(13)
                para.font.color.rgb = DARK
                para.space_after = Pt(6)

            else:  # paragraph
                para = body_frame.add_paragraph() if j > 0 else body_frame.paragraphs[0]
                para.text = content
                para.font.size = Pt(14)
                para.font.color.rgb = DARK
                para.space_after = Pt(6)

        # 페이지 번호
        pn_box = slide.shapes.add_textbox(
            Inches(12.0), Inches(7.0), Inches(1.0), Inches(0.4)
        )
        pn_frame = pn_box.text_frame
        pn_para = pn_frame.paragraphs[0]
        pn_para.text = str(i + 1)
        pn_para.font.size = Pt(10)
        pn_para.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
        pn_para.alignment = PP_ALIGN.RIGHT

    # 저장
    prs.save(str(PPTX_OUTPUT))
    size_mb = PPTX_OUTPUT.stat().st_size / 1024 / 1024
    print(f"   ✅ PPTX 저장: {PPTX_OUTPUT} ({size_mb:.1f}MB)")

    return PPTX_OUTPUT


# ─── 메인 ──────────────────────────────────────────────────────

def main():
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    # OCR/번역 캐시 확인 (중간에 중단됐을 때 재시작 지원)
    ocr_cache = OUTPUT_DIR / "ocr_results.json"
    tr_cache = OUTPUT_DIR / "translated_results.json"

    # 1단계: PDF → 이미지
    img_dir = OUTPUT_DIR / "pages"
    existing_pages = sorted(img_dir.glob("page_*.png")) if img_dir.exists() else []

    if existing_pages:
        print(f"📄 기존 이미지 {len(existing_pages)}페이지 발견, 변환 건너뜀")
        img_paths = existing_pages
    else:
        img_paths = pdf_to_images()

    # 2단계: OCR
    if ocr_cache.exists():
        print(f"\n🔍 기존 OCR 결과 발견, 로드 중...")
        ocr_results = json.loads(ocr_cache.read_text(encoding="utf-8"))
        print(f"   {len(ocr_results)}페이지 OCR 캐시 로드 완료")
    else:
        ocr_results = ocr_pages(img_paths)

    # 3단계: 번역
    if tr_cache.exists():
        print(f"\n🌐 기존 번역 결과 발견, 로드 중...")
        translated = json.loads(tr_cache.read_text(encoding="utf-8"))
        print(f"   {len(translated)}페이지 번역 캐시 로드 완료")
    else:
        translated = translate_pages(ocr_results)

    # 4단계: PPTX 생성
    pptx_path = create_pptx(translated, img_paths)

    print("\n" + "=" * 60)
    print("🎉 완료!")
    print(f"   📊 영문 PPTX: {pptx_path}")
    print(f"   📁 작업 폴더: {OUTPUT_DIR}")
    print("=" * 60)
    print("\n다음 단계:")
    print("  1. PPTX를 열어 내용/레이아웃 확인")
    print("  2. 고유명사(인명, 기관명) 수동 검수")
    print("  3. 필요 시 디자인 수정 후 PDF로 내보내기")


if __name__ == "__main__":
    main()
